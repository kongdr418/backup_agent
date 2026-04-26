"""
PPT 预览功能
将 PPT 转换为图片预览
使用 LibreOffice + Poppler 方案
"""

import os
import subprocess
import tempfile
import shutil
import base64
from typing import List, Dict
from pathlib import Path


class PPTPreviewer:
    """PPT 预览生成器 - 使用 LibreOffice + Poppler"""
    
    # 工具路径配置
    LIBREOFFICE_PATH = r"C:\Program Files\LibreOffice\program\soffice.exe"
    PDFTOPPM_PATH = r"D:\poppler\Release-25.12.0-0\poppler-25.12.0\Library\bin\pdftoppm.exe"
    
    def __init__(self, output_dir: str = "ppt_previews"):
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)
    
    def _check_tools(self) -> tuple[bool, str]:
        """检查必要工具是否可用"""
        if not os.path.exists(self.LIBREOFFICE_PATH):
            return False, f"LibreOffice 未找到: {self.LIBREOFFICE_PATH}"
        if not os.path.exists(self.PDFTOPPM_PATH):
            return False, f"Poppler (pdftoppm) 未找到: {self.PDFTOPPM_PATH}"
        return True, "工具检查通过"
    
    def generate_preview(self, pptx_path: str, max_slides: int = None) -> Dict:
        """
        生成 PPT 预览图片
        
        Args:
            pptx_path: PPT 文件路径
            max_slides: 最大预览页数，None 表示全部
            
        Returns:
            包含预览信息的字典
        """
        # 跳过 Office 临时文件（以 ~$ 开头）
        if os.path.basename(pptx_path).startswith('~$'):
            return {'error': '跳过临时文件', 'pptx_path': pptx_path}
        
        # 检查工具
        tools_ok, tools_msg = self._check_tools()
        if not tools_ok:
            print(f"⚠️ {tools_msg}")
            return {'error': tools_msg}
        
        if not os.path.exists(pptx_path):
            raise FileNotFoundError(f"PPT 文件不存在: {pptx_path}")
        
        # 创建临时目录
        temp_dir = tempfile.mkdtemp(prefix="ppt_preview_")
        
        try:
            # 1. PPT 转 PDF
            print(f"📄 正在转换 PPT 到 PDF...")
            pdf_path = self._pptx_to_pdf(pptx_path, temp_dir)
            
            if not pdf_path or not os.path.exists(pdf_path):
                return {'error': 'PPT 转 PDF 失败'}
            
            # 2. PDF 转图片
            print(f"🖼️ 正在生成预览图片...")
            base_name = Path(pptx_path).stem
            preview_dir = os.path.join(self.output_dir, base_name)
            
            # 清理旧文件
            if os.path.exists(preview_dir):
                shutil.rmtree(preview_dir)
            os.makedirs(preview_dir, exist_ok=True)
            
            images = self._pdf_to_images(pdf_path, preview_dir, max_slides)
            
            # 3. 转换为 base64 格式
            slide_images = []
            for i, img_path in enumerate(images, 1):
                with open(img_path, 'rb') as f:
                    img_base64 = base64.b64encode(f.read()).decode('utf-8')
                
                slide_images.append({
                    'page': i,
                    'path': img_path,
                    'base64': f"data:image/png;base64,{img_base64}",
                    'title': f'第 {i} 页'
                })
            
            return {
                'filename': Path(pptx_path).name,
                'total_pages': len(images),
                'slides': slide_images,
                'preview_dir': preview_dir
            }
            
        except Exception as e:
            print(f"预览生成失败: {e}")
            import traceback
            traceback.print_exc()
            return {'error': str(e)}
        
        finally:
            # 清理临时目录
            if os.path.exists(temp_dir):
                shutil.rmtree(temp_dir, ignore_errors=True)
    
    def _pptx_to_pdf(self, pptx_path: str, output_dir: str) -> str:
        """
        使用 LibreOffice 将 PPT 转换为 PDF
        
        Args:
            pptx_path: PPT 文件路径
            output_dir: 输出目录
            
        Returns:
            生成的 PDF 文件路径
        """
        cmd = [
            self.LIBREOFFICE_PATH,
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', output_dir,
            pptx_path
        ]
        
        try:
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=120
            )
            
            if result.returncode != 0:
                print(f"LibreOffice 错误: {result.stderr}")
                return None
            
            # 获取生成的 PDF 文件名
            base_name = Path(pptx_path).stem
            pdf_path = os.path.join(output_dir, f"{base_name}.pdf")
            
            return pdf_path if os.path.exists(pdf_path) else None
            
        except subprocess.TimeoutExpired:
            print("LibreOffice 转换超时")
            return None
        except Exception as e:
            print(f"LibreOffice 转换失败: {e}")
            return None
    
    def _pdf_to_images(self, pdf_path: str, output_dir: str, max_pages: int = None) -> List[str]:
        """
        使用 Poppler 将 PDF 转换为图片
        
        Args:
            pdf_path: PDF 文件路径
            output_dir: 输出目录
            max_pages: 最大页数，None 表示全部
            
        Returns:
            生成的图片路径列表
        """
        # 先生成所有图片
        output_prefix = os.path.join(output_dir, "slide")
        
        cmd = [
            self.PDFTOPPM_PATH,
            '-png',
            '-r', '150',  # 分辨率 150 DPI
            pdf_path,
            output_prefix
        ]
        
        try:
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=60
            )
            
            if result.returncode != 0:
                print(f"pdftoppm 错误: {result.stderr}")
                return []
            
            # 获取生成的图片文件
            images = sorted([
                os.path.join(output_dir, f)
                for f in os.listdir(output_dir)
                if f.startswith('slide-') and f.endswith('.png')
            ])
            
            # 重命名为更简洁的名称
            renamed_images = []
            for i, img_path in enumerate(images, 1):
                if max_pages and i > max_pages:
                    os.remove(img_path)
                    continue
                new_name = os.path.join(output_dir, f"slide_{i:02d}.png")
                os.rename(img_path, new_name)
                renamed_images.append(new_name)
            
            return renamed_images
            
        except subprocess.TimeoutExpired:
            print("pdftoppm 转换超时")
            return []
        except Exception as e:
            print(f"pdftoppm 转换失败: {e}")
            return []
    
    def get_preview_data(self, pptx_path: str) -> Dict:
        """
        获取预览数据（用于前端展示）
        
        Args:
            pptx_path: PPT 文件路径
            
        Returns:
            包含 base64 图片数据的字典
        """
        return self.generate_preview(pptx_path)
    
    def list_previews(self) -> List[str]:
        """列出所有可用的预览"""
        previews = []
        if os.path.exists(self.output_dir):
            for item in os.listdir(self.output_dir):
                item_path = os.path.join(self.output_dir, item)
                if os.path.isdir(item_path):
                    images = [f for f in os.listdir(item_path) if f.endswith('.png')]
                    if images:
                        previews.append(item)
        return previews


# 简单的文本预览（不需要外部工具）
def generate_text_preview(pptx_path: str) -> str:
    """
    生成 PPT 的文本预览（提取文字内容）
    
    Args:
        pptx_path: PPT 文件路径
        
    Returns:
        文本预览字符串
    """
    try:
        from pptx import Presentation
        
        prs = Presentation(pptx_path)
        preview_lines = [f"PPT: {Path(pptx_path).name}", f"共 {len(prs.slides)} 页", "=" * 50]
        
        for i, slide in enumerate(prs.slides, 1):
            preview_lines.append(f"\n【第 {i} 页】")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    preview_lines.append(f"  {shape.text.strip()}")
        
        return '\n'.join(preview_lines)
        
    except Exception as e:
        return f"文本预览生成失败: {e}"


if __name__ == "__main__":
    # 测试代码
    import sys
    import io
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
    
    previewer = PPTPreviewer()
    
    # 检查工具
    print("=" * 60)
    print("检查工具...")
    ok, msg = previewer._check_tools()
    print(f"{'✅' if ok else '❌'} {msg}")
    print("=" * 60)
    
    if ok:
        # 列出所有生成的 PPT
        ppt_dir = "generated_ppt"
        if os.path.exists(ppt_dir):
            # 过滤掉临时文件
            ppt_files = [f for f in os.listdir(ppt_dir) 
                        if f.endswith('.pptx') and not f.startswith('~$')]
            
            if ppt_files:
                test_ppt = os.path.join(ppt_dir, ppt_files[0])
                print(f"\n测试预览: {test_ppt}")
                
                # 文本预览
                print("\n" + "=" * 60)
                print("文本预览:")
                print("=" * 60)
                preview_text = generate_text_preview(test_ppt)
                # 保存到文件
                with open('preview.txt', 'w', encoding='utf-8') as f:
                    f.write(preview_text)
                print(preview_text[:1000] + "..." if len(preview_text) > 1000 else preview_text)
                print("\n(完整预览内容已保存到 preview.txt)")
                
                # 图片预览
                print("\n" + "=" * 60)
                print("生成图片预览...")
                result = previewer.generate_preview(test_ppt)
                
                if 'error' in result:
                    print(f"❌ 错误: {result['error']}")
                elif result.get('slides'):
                    print(f"✅ 生成成功: {len(result['slides'])} 张图片")
                    for slide in result['slides']:
                        print(f"  - 第 {slide['page']} 页: {slide['path']}")
                else:
                    print("❌ 图片预览生成失败")
            else:
                print("没有找到 PPT 文件")
        else:
            print("generated_ppt 目录不存在")
